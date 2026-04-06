import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─── Color Palette ───
DARK_BG       = RGBColor(11, 14, 20)
CARD_BG       = RGBColor(22, 27, 38)
INDIGO        = RGBColor(79, 70, 229)
INDIGO_LIGHT  = RGBColor(99, 102, 241)
WHITE         = RGBColor(240, 240, 245)
GRAY          = RGBColor(160, 170, 185)
ACCENT_CYAN   = RGBColor(56, 189, 248)
ACCENT_GREEN  = RGBColor(52, 211, 153)
ACCENT_AMBER  = RGBColor(251, 191, 36)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ─── Helpers ───
def dark_bg(slide, color=DARK_BG):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=WHITE,
             bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=WHITE, icon_color=INDIGO_LIGHT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        # Bullet icon
        run_icon = p.add_run()
        run_icon.text = "▸  "
        run_icon.font.size = Pt(font_size)
        run_icon.font.color.rgb = icon_color
        run_icon.font.name = font_name
        # Text
        run_text = p.add_run()
        run_text.text = item
        run_text.font.size = Pt(font_size)
        run_text.font.color.rgb = color
        run_text.font.name = font_name
    return txBox

def add_accent_line(slide, left, top, width, color=INDIGO):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_card(slide, left, top, width, height, title, items, icon_text="",
             accent=INDIGO_LIGHT):
    card = add_rect(slide, left, top, width, height, CARD_BG)
    add_accent_line(slide, left + Inches(0.3), top + Inches(0.15), Inches(0.8), accent)
    if icon_text:
        add_text(slide, left + Inches(0.3), top + Inches(0.3), Inches(0.5), Inches(0.5),
                 icon_text, font_size=22, color=accent, bold=True)
    add_text(slide, left + Inches(0.3) + (Inches(0.5) if icon_text else 0),
             top + Inches(0.3), width - Inches(1), Inches(0.5),
             title, font_size=18, color=WHITE, bold=True)
    add_bullet_list(slide, left + Inches(0.3), top + Inches(0.85),
                    width - Inches(0.6), height - Inches(1.2),
                    items, font_size=13, color=GRAY, icon_color=accent)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 – TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
dark_bg(slide)

# Decorative gradient rectangle at top
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

# Title
add_text(slide, Inches(0.8), Inches(1.2), Inches(7), Inches(1.2),
         "🚀 Orion-X", font_size=54, color=INDIGO_LIGHT, bold=True)

# Subtitle
add_text(slide, Inches(0.8), Inches(2.5), Inches(7), Inches(1),
         "The AI-Powered Learning Ecosystem", font_size=30, color=WHITE, bold=True)

# Tagline
add_text(slide, Inches(0.8), Inches(3.5), Inches(7), Inches(0.8),
         "Transforming static PDFs into dynamic, personalized\nlearning experiences using Generative AI.",
         font_size=16, color=GRAY)

# Tech badges
badges = "FastAPI  •  React 19  •  Google Gemini  •  PostgreSQL  •  Cloudinary"
add_text(slide, Inches(0.8), Inches(4.8), Inches(7), Inches(0.5),
         badges, font_size=12, color=INDIGO_LIGHT, bold=True)

# Hero image
hero_path = r"C:\Users\HP\.gemini\antigravity\brain\46ac0ff2-4c06-4899-a793-1f7efedc1611\orion_x_presentation_hero_1774951630390.png"
if os.path.exists(hero_path):
    slide.shapes.add_picture(hero_path, Inches(8.2), Inches(1), width=Inches(4.5))

# Bottom accent
add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 – THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(8), Inches(0.8),
         "🛑  The Problem", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(2))

problems = [
    ("📄", "Static Content", "PDFs sit passively — students read, forget, repeat.", ACCENT_CYAN),
    ("👥", "One-Size-Fits-All", "Every student gets the same experience regardless of pace.", ACCENT_GREEN),
    ("⏰", "Teacher Overload", "Manually creating quizzes and study guides is exhausting.", ACCENT_AMBER),
    ("😴", "Low Engagement", "Students disengage from dense, uninteractive material.", INDIGO_LIGHT),
]

for i, (icon, title, desc, accent) in enumerate(problems):
    col = i % 2
    row = i // 2
    left = Inches(0.8) + col * Inches(6.2)
    top  = Inches(1.8) + row * Inches(2.6)

    add_rect(slide, left, top, Inches(5.8), Inches(2.2), CARD_BG)
    add_accent_line(slide, left + Inches(0.3), top + Inches(0.2), Inches(0.8), accent)
    add_text(slide, left + Inches(0.3), top + Inches(0.4), Inches(0.6), Inches(0.6),
             icon, font_size=28, color=accent)
    add_text(slide, left + Inches(0.9), top + Inches(0.4), Inches(4.5), Inches(0.5),
             title, font_size=20, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.2), Inches(5.2), Inches(0.8),
             desc, font_size=14, color=GRAY)

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 – THE SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
         "✅  The Solution — Orion-X", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(2.5))

add_text(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.8),
         "An AI-native Learning Management System that doesn't just host content — it understands it.",
         font_size=18, color=GRAY)

solutions = [
    ("🧠", "Dynamic Content Ingestion", "PDFs are parsed, indexed, and turned\ninto searchable knowledge.", ACCENT_CYAN),
    ("🎯", "Personalized Study Paths", "AI generates custom roadmaps with\nsummaries & time estimates.", ACCENT_GREEN),
    ("💬", "Interactive AI Tutor", "RAG-based chatbot that answers\nquestions from YOUR notes.", INDIGO_LIGHT),
    ("⚡", "Instant Assessments", "Generate quizzes in seconds,\nnot hours.", ACCENT_AMBER),
]

for i, (icon, title, desc, accent) in enumerate(solutions):
    left = Inches(0.5) + i * Inches(3.15)
    top  = Inches(2.8)
    add_rect(slide, left, top, Inches(2.95), Inches(4), CARD_BG)
    add_accent_line(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), accent)
    add_text(slide, left + Inches(0.3), top + Inches(0.4), Inches(0.6), Inches(0.6),
             icon, font_size=30, color=accent)
    add_text(slide, left + Inches(0.3), top + Inches(1.1), Inches(2.3), Inches(0.6),
             title, font_size=16, color=WHITE, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.9), Inches(2.3), Inches(1.5),
             desc, font_size=13, color=GRAY)

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 – KEY FEATURES: TEACHERS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
         "👨‍🏫  Key Features — For Teachers", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(3))

add_card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.4),
         "📄  Smart Content Ingestion", [
             "Upload PDFs directly to Cloudinary",
             "Auto-extract text with PyMuPDF",
             "Seamless content pipeline"
         ], accent=ACCENT_CYAN)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.4),
         "🤖  AI Quiz Architect", [
             "MCQ generation via Gemini 2.0 Flash",
             "Based on actual uploaded notes",
             "Saves hours of manual work"
         ], accent=ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.4),
         "📹  Virtual Classrooms", [
             "Instant Jitsi Meet integration",
             "One-click meeting launch",
             "No external tools needed"
         ], accent=ACCENT_AMBER)

add_card(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.4),
         "📊  Performance Analytics", [
             "Automated student scoring",
             "Real-time progress tracking",
             "Actionable insights"
         ], accent=INDIGO_LIGHT)

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 – KEY FEATURES: STUDENTS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
         "🎓  Key Features — For Students", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(3))

add_card(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.4),
         "💬  Interactive AI Tutor (RAG)", [
             "Chat with AI about YOUR course materials",
             "Retrieval-Augmented Generation",
             "Context-aware, accurate responses"
         ], accent=ACCENT_CYAN)

add_card(slide, Inches(6.8), Inches(1.8), Inches(5.8), Inches(2.4),
         "🗺️  Personalized Study Paths", [
             "AI-generated learning roadmaps",
             "Summaries, key terms, study estimates",
             "Tailored to each module"
         ], accent=ACCENT_GREEN)

add_card(slide, Inches(0.5), Inches(4.5), Inches(5.8), Inches(2.4),
         "⏱️  Real-time Quizzing", [
             "Timed quizzes with instant feedback",
             "Track your progress over time",
             "Challenge yourself continuously"
         ], accent=ACCENT_AMBER)

add_card(slide, Inches(6.8), Inches(4.5), Inches(5.8), Inches(2.4),
         "✨  Premium UI/UX", [
             "Modern dark-mode design",
             "Framer Motion animations",
             "Responsive across devices"
         ], accent=INDIGO_LIGHT)

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 – TECHNICAL STACK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
         "🛠️  Technical Stack", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(2))

stack_items = [
    ("Frontend", "React 19, Vite, Framer Motion\nAxios, Lucide React Icons", ACCENT_CYAN),
    ("Backend", "FastAPI (Python 3.12)\nSQLAlchemy ORM, JWT Auth", ACCENT_GREEN),
    ("AI / ML", "Google Gemini Pro\nRAG Implementation", INDIGO_LIGHT),
    ("Database", "PostgreSQL / Supabase\nRelational + Cloud-hosted", ACCENT_AMBER),
    ("Storage", "Cloudinary\nPDF & Media Hosting", ACCENT_CYAN),
    ("Video", "Jitsi Meet Integration\nReal-time Conferencing", ACCENT_GREEN),
]

for i, (title, desc, accent) in enumerate(stack_items):
    col = i % 3
    row = i // 3
    left = Inches(0.5) + col * Inches(4.2)
    top  = Inches(1.8) + row * Inches(2.7)
    add_rect(slide, left, top, Inches(3.9), Inches(2.3), CARD_BG)
    add_accent_line(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), accent)
    add_text(slide, left + Inches(0.3), top + Inches(0.4), Inches(3.3), Inches(0.5),
             title, font_size=20, color=accent, bold=True)
    add_text(slide, left + Inches(0.3), top + Inches(1.0), Inches(3.3), Inches(1.0),
             desc, font_size=14, color=GRAY)

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 – ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.8),
         "🧠  System Architecture", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(2.5))

layers = [
    ("Client Layer", "React 19 SPA\nFramer Motion, Axios\nJWT Token Management", ACCENT_CYAN, Inches(0.5)),
    ("API Layer", "FastAPI Server\nBusiness Logic & Auth\nAI Orchestration (Gemini)", ACCENT_GREEN, Inches(4.5)),
    ("Data Layer", "PostgreSQL / Supabase\nSQLAlchemy ORM\nCloudinary Storage", ACCENT_AMBER, Inches(8.5)),
]

for title, desc, accent, left in layers:
    add_rect(slide, left, Inches(2), Inches(3.8), Inches(4.5), CARD_BG)
    add_accent_line(slide, left + Inches(0.2), Inches(2.15), Inches(0.6), accent)
    add_text(slide, left + Inches(0.3), Inches(2.4), Inches(3.2), Inches(0.5),
             title, font_size=22, color=accent, bold=True)
    add_text(slide, left + Inches(0.3), Inches(3.2), Inches(3.2), Inches(2.8),
             desc, font_size=14, color=GRAY)

# Arrows between layers
for x_start in [Inches(4.4), Inches(8.4)]:
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x_start, Inches(3.8), Inches(0.5), Inches(0.4))
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = INDIGO_LIGHT
    arrow.line.fill.background()

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 – AI TUTOR DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0.8), Inches(0.5), Inches(6), Inches(0.8),
         "💬  AI Tutor — RAG in Action", font_size=36, color=INDIGO_LIGHT, bold=True)
add_accent_line(slide, Inches(0.8), Inches(1.25), Inches(2.5))

add_text(slide, Inches(0.8), Inches(1.6), Inches(5.5), Inches(1),
         "The AI Tutor uses Retrieval-Augmented Generation to provide\naccurate, context-aware answers from your exact course PDFs.",
         font_size=16, color=GRAY)

add_bullet_list(slide, Inches(0.8), Inches(2.8), Inches(5.5), Inches(3), [
    "Student uploads question",
    "System retrieves relevant PDF chunks",
    "Gemini generates contextual answer",
    "Response cites source material",
    "Supports follow-up conversations"
], font_size=15, color=WHITE, icon_color=ACCENT_CYAN)

tutor_path = r"C:\Users\HP\.gemini\antigravity\brain\46ac0ff2-4c06-4899-a793-1f7efedc1611\orion_x_ai_tutor_mockup_1774951771346.png"
if os.path.exists(tutor_path):
    slide.shapes.add_picture(tutor_path, Inches(7), Inches(1.5), width=Inches(5.5))

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 – THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
dark_bg(slide)
add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), INDIGO)

add_text(slide, Inches(0), Inches(1.5), SLIDE_W, Inches(1.2),
         "🚀 Thank You!", font_size=54, color=INDIGO_LIGHT, bold=True,
         alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(3), SLIDE_W, Inches(0.8),
         "Orion-X: The AI-Powered Learning Ecosystem", font_size=24, color=WHITE,
         bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(0), Inches(4), SLIDE_W, Inches(0.8),
         '"The transition from hosting information to providing intelligence."',
         font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(4.5), Inches(5.2), Inches(4.3), Inches(1.2), CARD_BG)
add_text(slide, Inches(4.5), Inches(5.3), Inches(4.3), Inches(1),
         "Questions & Discussion", font_size=22, color=ACCENT_CYAN, bold=True,
         alignment=PP_ALIGN.CENTER)

add_rect(slide, Inches(0), Inches(7.42), SLIDE_W, Inches(0.08), INDIGO)

# ─── Save ───
output_path = r"c:\Users\HP\Desktop\Projects\orion-x\Orion-X_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved: {output_path}")
print(f"   Slides: {len(prs.slides)}")
